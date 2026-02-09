from dataclasses import dataclass
import sys
import subprocess
import shlex
from pathlib import Path
from typing import Iterable
import re
from zipfile import ZipFile, ZIP_DEFLATED
import os
from io import BytesIO
from copy import deepcopy
from uuid import uuid4

from lxml import etree
from lxml.builder import ElementMaker

import pptx  # python-pptx
import pptx.util
import pptx.shapes
from pptx.shapes.picture import Picture
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from slugify import slugify

from bakery_scripts.utils import get_mime_type, patch_math_for_pandoc
from bakery_scripts import mathml2png

import imgkit

from PIL import Image

from . import excepthook


excepthook.attach(sys)


NS_XHTML = "http://www.w3.org/1999/xhtml"
E = ElementMaker(namespace=NS_XHTML, nsmap={None: NS_XHTML})
BLOCKISH_TAGS = (
    f"{{{NS_XHTML}}}address",
    f"{{{NS_XHTML}}}article",
    f"{{{NS_XHTML}}}aside",
    f"{{{NS_XHTML}}}blockquote",
    f"{{{NS_XHTML}}}details",
    f"{{{NS_XHTML}}}dialog",
    f"{{{NS_XHTML}}}dd",
    f"{{{NS_XHTML}}}div",
    f"{{{NS_XHTML}}}dl",
    f"{{{NS_XHTML}}}dt",
    f"{{{NS_XHTML}}}fieldset",
    f"{{{NS_XHTML}}}figcaption",
    f"{{{NS_XHTML}}}figure",
    f"{{{NS_XHTML}}}footer",
    f"{{{NS_XHTML}}}form",
    f"{{{NS_XHTML}}}h1",
    f"{{{NS_XHTML}}}h2",
    f"{{{NS_XHTML}}}h3",
    f"{{{NS_XHTML}}}h4",
    f"{{{NS_XHTML}}}h5",
    f"{{{NS_XHTML}}}h6",
    f"{{{NS_XHTML}}}header",
    f"{{{NS_XHTML}}}hgroup",
    f"{{{NS_XHTML}}}hr",
    f"{{{NS_XHTML}}}li",
    f"{{{NS_XHTML}}}main",
    f"{{{NS_XHTML}}}nav",
    f"{{{NS_XHTML}}}ol",
    f"{{{NS_XHTML}}}p",
    f"{{{NS_XHTML}}}pre",
    f"{{{NS_XHTML}}}section",
    f"{{{NS_XHTML}}}table",
    f"{{{NS_XHTML}}}tr",
    f"{{{NS_XHTML}}}td"
    f"{{{NS_XHTML}}}ul",
)


def info(msg):
    sys.stderr.write(msg)
    sys.stderr.flush()


def replace_blockish(elem: etree.ElementBase):
    blockish_elements = (e for e in elem.iter(None) if e.tag in BLOCKISH_TAGS)
    for blockish in blockish_elements:
        blockish.tag = f"{{{NS_XHTML}}}span"
    return elem


def to_slide_safe(content: str | etree.ElementBase):
    if isinstance(content, str):
        return content
    return replace_blockish(content)


def class_xpath(class_name: str):
    return f'contains(concat(" ", @class, " "), " {class_name} ")'


def try_find_nearest_sm(elem):
    sm_results = elem.xpath("ancestor-or-self::*[@data-sm][1]/@data-sm")
    if not sm_results:
        return "N/A"
    return sm_results[0]


def get_elem_text(cell):
    return "".join(cell.itertext()).strip()


class Element:
    def __init__(self, element):
        self._element = element

    def has_class(self, class_name: str):
        class_list = self.element.get("class", "").strip()
        return f" {class_name} " in f" {class_list} "

    @property
    def element(self):
        return self._element

    def get_doc_dir(self):
        return Path(self.element.getroottree().docinfo.URL).parent

    def xpath(self, p: str, *, namespaces=None):
        _namespaces = {"h": "http://www.w3.org/1999/xhtml"}
        if namespaces is not None:  # pragma: no cover
            _namespaces.update(namespaces)
        return self._element.xpath(p, namespaces=_namespaces)

    def xpath1(self, p: str):
        return self.xpath(p)[0]

    def xpath1_or_none(self, p: str):
        try:
            return self.xpath1(p)
        except IndexError:
            return None


class BookElement(Element):
    def get_title(self):
        return self.xpath1(
            './/*[@data-type="metadata"]/*[@data-type="document-title"]/text()'
        ).strip()

    def get_number(self):  # pragma: no cover
        raise NotImplementedError()


class Captioned(Element):
    def get_caption_elem(self):
        return self.xpath1(f'./*[{class_xpath("os-caption-container")}]')

    def has_caption(self):
        return (
            self.xpath1_or_none(f'.//*[{class_xpath("os-caption-container")}]')
            is not None
        )

    def has_number(self):
        return self.has_caption() and self.get_number() != ""

    def get_number(self) -> str:
        caption_number = self.get_caption_elem().xpath(
            f'./*[{class_xpath("os-number")}]//text()'
        )
        return "".join(caption_number).strip()

    def get_caption(self) -> str:
        caption_text = self.get_caption_elem().xpath(
            f'./*[{class_xpath("os-caption")}]//text()'
        )
        return "".join(caption_text).strip()

    def get_title_elem(self, prefix):
        title_query = self.get_caption_elem().xpath(
            f'./*[{class_xpath("os-title")}]'
        )
        number = self.get_number()
        span = E.span()
        span.text = " ".join(part for part in (prefix, number) if part)
        if title_query:
            span.text += " "
            span.append(title_query[0])
        return span


class Figure(Captioned):
    def get_img(self):
        return self.xpath1(".//h:img")

    def get_src(self) -> str | None:
        return self.get_img().get("src")

    def get_alt(self) -> str | None:
        return self.get_img().get("alt")


class Table(Captioned):
    def get_title(self):
        return super().get_title_elem("Table")

    def get_table_elem(self):
        return self.xpath1(".//h:table")

    def get_alt_text(self, slide_title) -> str:
        base_caption = self.get_caption()
        if base_caption:
            return base_caption

        # Generate caption from table structure when base_caption is empty
        table_elem = Element(self.get_table_elem())

        caption_parts = [
            get_elem_text(slide_title)
            if not isinstance(slide_title, str)
            else slide_title
        ]

        # Try to find headers in thead
        thead_rows = table_elem.xpath(".//h:thead//h:tr")
        if thead_rows:
            if len(thead_rows) > 1:
                for i, thead_row in enumerate(thead_rows, start=1):
                    header_cells = Element(thead_row).xpath(".//h:th|.//h:td")
                    if header_cells:
                        headers = [get_elem_text(cell) for cell in header_cells]
                        caption_parts.append(
                            f"Table heading row {i} columns: {', '.join(headers)}"
                        )
            else:
                thead_row = thead_rows[0]
                header_cells = Element(thead_row).xpath(".//h:th|.//h:td")
                if header_cells:
                    headers = [get_elem_text(cell) for cell in header_cells]
                    caption_parts.append(f"Columns: {', '.join(headers)}")

            data_rows = table_elem.xpath(".//h:tbody//h:tr")
        else:
            data_rows = table_elem.xpath(".//h:tr")

        # Extract data from rows
        for i, row in enumerate(data_rows, start=1):
            cells = Element(row).xpath(".//h:td|.//h:th")
            if cells:
                row_data = [get_elem_text(cell) for cell in cells]
                caption_parts.append(f"Row {i}: {', '.join(row_data)}")

        if len(caption_parts) > 1:
            return "; ".join(caption_parts)
        return "Table without data"  # pragma: no cover


class Page(BookElement):
    def __init__(self, parent_chapter: "Chapter", number: int, elem):
        super().__init__(elem)
        self.parent_chapter = parent_chapter
        self._number_in_chapter = number

    @property
    def number_less_intro(self):
        if self.parent_chapter.has_introduction:
            return self._number_in_chapter - 1
        return self._number_in_chapter

    def try_find_lo_container(self):
        sections = self.xpath(".//h:section")
        if sections:
            section = Element(sections[0])
            if (
                # Ideal case
                section.has_class("learning-objectives") or
                # Last chance: check title
                "".join(
                    section.xpath('.//*[@data-type = "title"][1]//text()')
                ).strip().lower() == "learning objectives"
            ):
                return section
        abstracts = self.xpath('.//*[@data-type = "abstract"]')
        if abstracts:
            return Element(abstracts[0])
        return None

    def get_learning_objectives(self):
        learning_objectives: list[str | etree.ElementBase] = []
        lo_container = self.try_find_lo_container()
        if lo_container is not None:
            lo_elements = lo_container.xpath(
                f'.//*[{class_xpath("os-abstract-content")}]'
            )
            if not lo_elements:
                lo_elements = lo_container.xpath(".//h:li")
            for lo_elem in lo_elements:
                span = E.span(*lo_elem)
                span.text = lo_elem.text
                learning_objectives.append(span)
        return learning_objectives

    def get_figures(self):
        figure_elems = self.xpath(".//h:figure/parent::*")
        return [Figure(elem) for elem in figure_elems]

    def get_number(self):
        return f"{self.parent_chapter.get_number()}.{self.number_less_intro}"

    def get_tables(self):
        # Get all tables that are not nested in other tables
        table_elems = self.xpath(
            ".//h:table[not(ancestor::h:table)]/parent::*"
        )
        return [Table(elem) for elem in table_elems]

    @property
    def is_introduction(self):
        return self.has_class("introduction")

    @property
    def is_summary(self):
        return self.get_title().lower().endswith("summary")


class Chapter(BookElement):
    def __init__(self, number, elem):
        super().__init__(elem)
        self.number = number
        self._pages = []
        self._has_introduction = None

    @property
    def has_introduction(self):
        if self._has_introduction is None:
            self._has_introduction = any(
                p.is_introduction for p in self.get_pages()
            )
        return self._has_introduction

    def get_number(self):
        return str(self.number)

    def get_pages(self):
        if not self._pages:
            page_elems = self.xpath('.//*[@data-type="page"]')
            pages = enumerate(page_elems, start=1)
            self._pages = [Page(self, i, elem) for i, elem in pages]
        return self._pages

    def get_chapter_outline(
        self, *, include_introduction=False, include_summary=False
    ):
        pages = self.get_pages()
        if not include_summary:
            pages = (page for page in pages if not page.is_summary)
        if not include_introduction:
            pages = (page for page in pages if not page.is_introduction)
        titles = (page.get_title().strip() for page in pages)
        non_empty_titles = (title for title in titles if title)
        return list(non_empty_titles)


class Book(Element):
    def get_chapters(self):
        chapters = enumerate(self.xpath('//*[@data-type="chapter"]'), start=1)
        return [Chapter(i, el) for i, el in chapters]

    def get_title(self):
        title_text = self.xpath(".//h:title//text()")
        return "".join(title_text)


def get_document_indices(elems: Iterable[Element]):
    elem_list = list(elems)
    tree = elem_list[0].element.getroottree()
    for idx, elem in enumerate(tree.iter()):
        found = None
        for elem_list_idx, search in enumerate(elem_list):
            if elem == search.element:
                yield idx, search
                found = elem_list_idx
                break
        if found is not None:
            del elem_list[found]
            if not elem_list:
                return


def sort_by_document_index(elems: Iterable[Element]):
    return (elem for _, elem in get_document_indices(elems))


@dataclass(kw_only=True)
class SlideContent:
    title: str | etree.ElementBase
    notes: str | None = None


@dataclass(kw_only=True)
class OutlineSlideContent(SlideContent):
    bullets: list[str | etree.ElementBase]
    heading: str | None = None
    numbered: bool = False
    number_offset: int = 1


@dataclass(kw_only=True)
class FigureSlideContent(SlideContent):
    src: str
    alt: str
    caption: str


@dataclass(kw_only=True)
class TableSlideContent(SlideContent):
    os_table: Table


@dataclass(kw_only=True)
class HTMLTableSlideContent(SlideContent):
    html: etree.ElementBase
    caption: str = ""


def guess_str_len_html(item: str | etree.ElementBase) -> int:
    if isinstance(item, str):
        content = item
    else:
        content = "".join(item.itertext(None))
    html_like_whitespace = re.sub(r"\s+", " ", content)
    return len(html_like_whitespace)


def chunk_bullets(
    bullets: list[str | etree.ElementBase],
    max_bullets: int,
    max_characters: int
):
    character_count = 0
    buffer = []
    length_by_bullet_idx = [guess_str_len_html(b) for b in bullets]
    avg_bullet_length = sum(length_by_bullet_idx) / len(bullets)
    for i, bullet in enumerate(bullets):
        bullet_len = length_by_bullet_idx[i]
        bullet_len *= bullet_len / avg_bullet_length
        if buffer and (
            len(buffer) >= max_bullets or
            character_count + bullet_len >= max_characters
        ):
            yield buffer
            buffer = [bullet]
            character_count = bullet_len
        else:
            buffer.append(bullet)
            character_count += bullet_len
    if buffer:
        yield buffer


def split_large_bullet_lists(slide_contents: Iterable[SlideContent]):
    for slide_content in slide_contents:
        if isinstance(slide_content, OutlineSlideContent):
            if slide_content.heading is not None:
                max_bullets, max_chars = 6, 300
            else:
                max_bullets, max_chars = 9, 400
            bullet_groups = list(
                chunk_bullets(slide_content.bullets, max_bullets, max_chars)
            )
            slide_count = len(bullet_groups)
            if slide_count > 1:
                number_offset = 1
                for i, bullets in enumerate(bullet_groups, start=1):
                    title = f"{slide_content.title} ({i} of {slide_count})"
                    yield OutlineSlideContent(
                        title=title,
                        heading=slide_content.heading,
                        bullets=bullets,
                        numbered=slide_content.numbered,
                        number_offset=number_offset,
                    )
                    number_offset += len(bullets)
                continue
        # Default to yielding original content
        yield slide_content


def index_to_coord(idx, image_width):
    coord_raw = idx / image_width
    row = int(coord_raw)
    col = round((coord_raw - row) * image_width)
    return col, row


def image_to_pixels(img: Image.Image):
    b = img.tobytes()
    stride = len(img.getpixel((0, 0)))
    pixels = zip(*(b[offset::stride] for offset in range(stride)))
    enumerated_pixels = enumerate(pixels)
    return enumerated_pixels


def auto_crop_img(img, bg=None):
    img_width = img.width
    pixels = image_to_pixels(img)
    if bg is None:
        bg = img.getpixel((0, 0))
    non_bg_coords = (
        index_to_coord(idx, img_width)
        for idx, _ in (p for p in pixels if p[1] != bg)
    )
    # Intentionally reversed so we can find the min/max
    right, bottom, left, top = img.getbbox()
    # Create largest box that contains any non-background pixels
    for col, row in non_bg_coords:
        if col < left:
            left = col
        elif col > right:
            right = col
        if row < top:
            top = row
        elif row > bottom:
            bottom = row
    box = (left, top, right + 1, bottom + 1)
    return img.crop(box)


# NOTE: Sometimes this can cause a 399 response from Google Fonts. This may be
# due to some kind of rate limiting. If this becomes a problem, a possible
# workaround is to download the font css once and update the references in the
# pdf css so that the font is read from the local copy
def xhtml_to_img(elem, *, css=[]):
    options = {
        "format": "png",
        "log-level": "warn",
        "quality": "100",
        "enable-local-file-access": "",
    }
    with BytesIO() as img_file:
        serialized = etree.tostring(elem, encoding="unicode")
        assert serialized, "Cannot convert empty element"
        img_bytes = imgkit.from_string(
            serialized, None, css=css, options=options
        )
        img_file.write(img_bytes)
        img_file.seek(0)
        img = Image.open(img_file)
        img.load()  # load eagerly (the BytesIO will close)
    return img


def element_to_image(
    elem: etree.ElementBase, doc_dir: Path, resource_dir: Path, css: list[str]
) -> Image.Image:
    math_nodes = elem.xpath(
        "//m:math", namespaces={"m": "http://www.w3.org/1998/Math/MathML"}
    )
    # NOTE: Requires mathjax json rpc to be running
    # This replaces mathml with imgs that have relative paths
    mathml2png.convert_math(math_nodes, resource_dir, use_svg=True)
    # NOTE: wkhtmltoimage requires absolute paths
    for resource in elem.xpath('.//*[@src]'):
        attr = "src"
        rel_p = resource.get(attr)
        abs_p = (doc_dir / rel_p).resolve()
        resource.set(attr, str(abs_p))
    img = xhtml_to_img(elem, css=css)
    img = auto_crop_img(img)
    return img


def os_table_to_image(
    os_table: Table, resource_dir: Path, css: list[str]
) -> Image.Image:
    doc_dir = os_table.get_doc_dir()
    # Clone the .os-table div. The div tends to be part of the css selector.
    table_clone = deepcopy(os_table.element)
    # Remove everything other than the table from the div (we do not want the
    # caption in the image)
    for elem in list(table_clone):
        if elem.tag != "{http://www.w3.org/1999/xhtml}table":
            table_clone.remove(elem)
    return element_to_image(table_clone, doc_dir, resource_dir, css)


def handle_tables(
    slide_contents: Iterable[SlideContent], resource_dir: Path, css: list[str]
):
    for slide_content in slide_contents:
        if isinstance(slide_content, TableSlideContent):
            table_slide = slide_content
            os_table = table_slide.os_table
            title = table_slide.title
            img = os_table_to_image(os_table, resource_dir, css)
            w, h = img.size
            if (
                h > 250 or
                w * h > 200000 or
                len(os_table.xpath(".//h:table")) > 1 or  # nested tables
                len(os_table.xpath(".//h:img|.//h:iframe|.//h:video")) > 0
            ):
                doc_dir = os_table.get_doc_dir()
                img_name = f"{uuid4()}.png"
                img_path = resource_dir / img_name
                img.save(img_path)

                note_text = None
                alt_text = os_table.get_alt_text(title)
                if len(alt_text) > 200:
                    ellipsis = "... (Full description in notes)"
                    note_text = alt_text
                    alt_text = alt_text[:200 - len(ellipsis)] + ellipsis

                yield FigureSlideContent(
                    title=title,
                    src=os.path.relpath(img_path, doc_dir),
                    caption=os_table.get_caption(),
                    alt=alt_text,
                    notes=note_text,
                )
            else:
                yield HTMLTableSlideContent(
                    title=title,
                    html=os_table.get_table_elem(),
                    caption=os_table.get_caption()
                )
        else:
            yield slide_content


def rename_images_to_type(slide_contents: Iterable[SlideContent], resource_dir):
    extension_by_mime_type = {
        "image/jpeg": ".jpg",
        "image/png": ".png",
        "image/gif": ".gif",
        "image/tiff": ".tiff",
        "image/svg+xml": ".svg",
    }
    for slide_content in slide_contents:
        if isinstance(slide_content, FigureSlideContent):
            fig = slide_content
            resource_name = Path(fig.src).name
            resource_src = Path(resource_dir) / resource_name
            assert resource_src.exists(), f"Missing resource: {resource_src}"
            if resource_src.suffix == "":
                mime_type = get_mime_type(str(resource_src))
                ext = extension_by_mime_type.get(mime_type, "")
                if ext != "":
                    resource_dst = resource_src.with_suffix(ext)
                    if not resource_dst.exists():
                        os.link(resource_src, resource_dst)
                    yield FigureSlideContent(
                        title=fig.title,
                        notes=fig.notes,
                        src=fig.src.replace(resource_name, resource_name + ext),
                        alt=fig.alt,
                        caption=fig.caption,
                    )
                    continue
        yield slide_content


def chapter_to_slide_contents(chapter: Chapter):
    if chapter.get_chapter_outline():
        yield OutlineSlideContent(
            title="Chapter outline",
            bullets=chapter.get_chapter_outline(),
            numbered=True,
        )
    for page in chapter.get_pages():
        info(".")
        learning_objectives = page.get_learning_objectives()
        figures = page.get_figures()
        tables = page.get_tables()
        page_contents = figures + tables
        if learning_objectives:
            title_parts = (
                p
                for p in (page.get_number(), page.get_title())
                if p is not None
            )
            if page.is_summary:  # pragma: no cover
                heading = None
            else:
                heading = "Learning Objectives"
            yield OutlineSlideContent(
                title=" ".join(title_parts),
                heading=heading,
                bullets=learning_objectives,
            )
        if page_contents:
            for page_content in sort_by_document_index(page_contents):
                if isinstance(page_content, Figure):
                    fig = page_content
                    if not fig.has_number():
                        continue
                    src = fig.get_src()
                    title = f"Figure {fig.get_number()}"
                    caption = fig.get_caption() or fig.get_alt() or "None"
                    alt = fig.get_alt() or ""
                    if not src:  # pragma: no cover
                        name = "src"
                        parent_page_id = fig.element.xpath(
                            'ancestor::*[@data-type = "page"]/@id'
                        )[0]
                        data_sm = try_find_nearest_sm(fig.get_img())
                        raise Exception(
                            f"Figure missing {name}\n"
                            f"Page ID: {parent_page_id}\n"
                            f"Nearest element location: {data_sm}"
                        )
                    yield FigureSlideContent(
                        title=title,
                        src=src,
                        caption=caption,
                        alt=alt,
                        notes=alt,
                    )
                elif isinstance(page_content, Table):
                    table = page_content
                    if not table.has_number():  # pragma: no cover
                        continue
                    yield TableSlideContent(
                        title=table.get_title(), os_table=table,
                    )


def slide_contents_to_html(
    title: str, subtitle: str, slide_contents: Iterable[SlideContent]
):
    def slide_content_to_html_parts(slide_contents: Iterable[SlideContent]):
        for slide_content in slide_contents:
            yield E.h2(to_slide_safe(slide_content.title))
            if isinstance(slide_content, (OutlineSlideContent,)):
                if slide_content.heading is not None:
                    yield E.strong(slide_content.heading)
                list_items = [
                    E.li(to_slide_safe(item)) for item in slide_content.bullets
                ]
                if slide_content.numbered:
                    yield E.ol(
                        *list_items, start=str(slide_content.number_offset)
                    )
                else:
                    yield E.ul(*list_items)
            elif isinstance(slide_content, FigureSlideContent):
                yield E.figure(
                    E.img(
                        src=slide_content.src,
                        alt=slide_content.caption,
                        title=slide_content.alt,
                    )
                )
            elif isinstance(slide_content, HTMLTableSlideContent):
                table_elem = slide_content.html
                if slide_content.caption:
                    table_caption = E.caption(slide_content.caption)
                    table_elem.insert(0, table_caption)
                yield table_elem
            else:  # pragma: no cover
                raise TypeError(
                    f"Cannot convert slide content type: {type(slide_content)}"
                )
            if slide_content.notes is not None:
                notes = slide_content.notes
                html_notes = [E.div(note) for note in notes.split("\n")]
                yield E.div(*html_notes, **{"class": "notes"})

    return E.html(
        E.head(E.title(title), E.meta(name="subtitle", content=subtitle)),
        E.body(*slide_content_to_html_parts(slide_contents)),
    )


def slides_etree_to_ppt(
    slides_etree, html_output, ppt_output, resource_dir, reference_doc
):  # pragma: no cover
    with open(html_output, "wb") as fout:
        slides_etree.getroottree().write(
            fout,
            encoding="utf-8",
            xml_declaration=False,
            pretty_print=True,
        )
    cmd = shlex.split(
        "pandoc "
        "--fail-if-warnings "
        "--from=html "
        "--to=pptx "
        f'--resource-path "{resource_dir}" '
        f'--reference-doc "{reference_doc}" '
        f'--output "{ppt_output}" '
        f'"{html_output}" '
    )
    subprocess.run(cmd, check=True, stderr=sys.stderr, stdout=sys.stdout)


def get_descr(shape: BaseShape):
    matches = shape.element.xpath("//*[@descr]")
    assert len(matches) == 1, f"Expected one descr, got {len(matches)}"
    return matches[0]


def get_alt_text(shape: BaseShape):
    elem = get_descr(shape)
    return elem.get("descr")


def set_alt_text(shape: BaseShape, alt_text: str):
    elem = get_descr(shape)
    elem.set("descr", alt_text)


def insert_cover_image(pres: pptx.Presentation, cover_image: str):
    height = pptx.util.Inches(3.5)
    pic = pres.slides[0].shapes.add_picture(
        cover_image,
        height=height,
        top=round(pres.slide_height * 5 / 7 - height / 2),
        left=0,
    )
    pic.left = round(pres.slide_width / 2 - pic.width / 2)
    pic.name = "Cover Image"
    set_alt_text(pic, "Cover image")


def fix_image_alt_text(slides: Iterable[Slide]):
    for slide in slides:
        for shape in slide.shapes:
            if isinstance(shape, Picture):
                alt_text = get_alt_text(shape)
                # Remove image path from alt text because otherwise powerpoint
                # does not recognize the alt text as existing
                alt_text_no_image_path = re.sub(
                    r"\s*\.{1,2}\/[^.]+\.[a-z]+$",
                    "",
                    alt_text,
                    flags=re.IGNORECASE,
                )
                set_alt_text(shape, alt_text_no_image_path)
        yield slide


def adjust_figure_caption_font(slides: Iterable[Slide]):
    for slide in slides:
        title = slide.shapes.title.text.lower()
        if title.startswith("figure") or title.startswith("table"):
            caption = slide.shapes[-1]
            if caption.has_text_frame:
                for p in caption.text_frame.paragraphs:
                    p.font.size = pptx.util.Pt(12)
        yield slide


def slides_post_process(pres: pptx.Presentation):
    slides = pres.slides
    slides = fix_image_alt_text(slides)
    slides = adjust_figure_caption_font(slides)
    _ = list(slides)  # Run generator to completion
    slide_layouts_by_name = {sl.name: sl for sl in pres.slide_layouts}
    last_slide_layout = slide_layouts_by_name["Last Slide"]
    pres.slides.add_slide(last_slide_layout)


def fix_namespaces(tree):
    # These are tag names mapped to namespaces that pandoc fails to add when
    # converting to pptx. Normally lxml would resolve namespace prefixes to the
    # namespace. Since these tags use undefined prefixes, they are not resolved
    ns_fixes = {
        "a14:m": etree.QName(
            "http://schemas.microsoft.com/office/drawing/2010/main", "m"
        ),
    }
    for elem in tree.iter():
        if elem.tag in ns_fixes:
            qname = ns_fixes[elem.tag]
            nsname = elem.tag.split(":")[0]
            new_elem = etree.Element(
                qname, nsmap={nsname: qname.namespace}, attrib=elem.attrib
            )
            for child in elem.iterchildren():
                new_elem.append(child)
            parent = elem.getparent()
            parent.replace(elem, new_elem)


def is_slide_filename(filename):
    return re.match(r".+slides/slide\d+\.xml$", filename) is not None


def fix_pptx_file(input_path: Path):  # pragma: no cover
    lax_xml_parser = etree.XMLParser(recover=True)
    output_path = input_path.with_suffix(".tmp")
    with (
        ZipFile(input_path) as pptx_i,
        ZipFile(output_path, mode="w", compression=ZIP_DEFLATED) as pptx_o,
    ):
        for info in pptx_i.infolist():
            with (
                pptx_i.open(info.filename, "r") as fin,
                pptx_o.open(info.filename, "w") as fout,
            ):
                content = fin.read()
                if is_slide_filename(info.filename):
                    xml = etree.XML(content, lax_xml_parser)
                    fix_namespaces(xml)
                    xml.getroottree().write(fout)
                else:
                    fout.write(content)
    output_path.rename(input_path)


def main():
    book_input = Path(sys.argv[1]).resolve(strict=True)
    resource_dir = Path(sys.argv[2]).resolve(strict=True)
    reference_doc = Path(sys.argv[3]).resolve(strict=True)
    cover_image = Path(sys.argv[4]).resolve()
    css = [sys.argv[5]]
    out_fmt = sys.argv[6]
    # For xhtml_to_img
    os.environ.setdefault("XDG_RUNTIME_DIR", "/tmp/runtime-root")
    tree = etree.parse(str(book_input), None)
    book = Book(tree.getroot())
    patch_math_for_pandoc(book.element, "http://www.w3.org/1998/Math/MathML")
    for chapter in book.get_chapters():
        title = f"{chapter.get_number()} {chapter.get_title()}".replace("'", "")
        slug = slugify(title)
        ppt_output, html_output = (
            Path(out_fmt.format(slug=slug, extension="pptx")),
            Path(out_fmt.format(slug=slug, extension="html")),
        )
        info(f"Working on: {slug}")
        slide_contents = chapter_to_slide_contents(chapter)
        slide_contents = split_large_bullet_lists(slide_contents)
        slide_contents = handle_tables(slide_contents, resource_dir, css)
        slide_contents = rename_images_to_type(slide_contents, resource_dir)
        slides_etree = slide_contents_to_html(
            book.get_title(),
            f"Chapter {chapter.get_number()} {chapter.get_title()}",
            slide_contents,
        )
        slides_etree_to_ppt(
            slides_etree, html_output, ppt_output, resource_dir, reference_doc
        )
        fix_pptx_file(ppt_output)
        tmp_path = Path(ppt_output).with_suffix(".pptx.tmp")
        pres = pptx.Presentation(ppt_output)
        slides_post_process(pres)
        if cover_image.exists():
            insert_cover_image(pres, str(cover_image))
        pres.save(str(tmp_path))
        tmp_path.rename(ppt_output)
        info("Done!\n")
